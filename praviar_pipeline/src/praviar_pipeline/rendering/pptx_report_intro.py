"""Introductory slides for PPTX report rendering."""

from __future__ import annotations

from typing import TYPE_CHECKING

from praviar_pipeline.rendering.brand_mark import render_praviar_mark_png_stream
from praviar_pipeline.rendering.branding import (
    SUPPORTED_OFFICE_BRANDING_LOGO_EXTENSIONS,
    resolve_branding_logo_path,
)
from praviar_pipeline.rendering.design import (
    BRAND_INK,
    BRAND_ON_INK_MUTED,
    BRAND_PAPER,
    BRAND_SECONDARY_TEXT,
    RISK_FILL,
    risk_display,
)
from praviar_pipeline.rendering.evidence_scope import (
    format_source_list,
    source_status_detail,
    source_status_label,
    summarize_evidence_scope,
)
from praviar_pipeline.rendering.governed_decision import (
    governed_blocking_count,
    governed_decision_label,
    governed_executive_summary,
    governed_risk_level,
)
from praviar_pipeline.rendering.pptx_report_layout import (
    add_ink_title_bar,
    add_speaker_notes,
    add_text_box,
    pptx_rgb,
    trunc,
)

if TYPE_CHECKING:
    from praviar_pipeline.models.report import FTOReport
    from praviar_pipeline.rendering.branding import BrandingConfig
    from praviar_pipeline.rendering.export_options import ExportRenderOptions


def add_cover_slide(prs, report: FTOReport, branding: BrandingConfig) -> None:
    """Slide 1: Navy cover with title, compound, date."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = pptx_rgb(branding.primary_color)
    bg.line.fill.background()

    logo_path = resolve_branding_logo_path(
        branding,
        renderer_name="PPTX",
        supported_extensions=SUPPORTED_OFFICE_BRANDING_LOGO_EXTENSIONS,
    )
    if logo_path is not None:
        slide.shapes.add_picture(
            str(logo_path),
            Inches(1),
            Inches(0.62),
            width=Inches(1.0),
        )
    elif not branding.suppresses_praviar_branding:
        slide.shapes.add_picture(
            render_praviar_mark_png_stream(variant="on_dark", size_px=384),
            Inches(1),
            Inches(0.58),
            width=Inches(1.0),
        )

    if not branding.suppresses_praviar_branding or branding.firm_name:
        add_text_box(
            slide,
            Inches(1.95),
            Inches(0.72),
            Inches(4.5),
            Inches(0.4),
            branding.display_name,
            font_size=16,
            bold=True,
            color_hex=BRAND_PAPER,
        )

    add_text_box(
        slide,
        Inches(1),
        Inches(1.5),
        Inches(11),
        Inches(1),
        "Freedom-to-Operate Analysis",
        font_size=36,
        bold=True,
        color_hex=BRAND_PAPER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(2.8),
        Inches(11),
        Inches(0.8),
        report.compound.name,
        font_size=24,
        color_hex=BRAND_PAPER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(4.0),
        Inches(11),
        Inches(0.5),
        f"Generated: {report.generated_at.strftime('%Y-%m-%d %H:%M UTC')}",
        font_size=14,
        color_hex=BRAND_ON_INK_MUTED,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(4.5),
        Inches(11),
        Inches(0.5),
        f"Report ID: {report.report_id}",
        font_size=11,
        color_hex=BRAND_ON_INK_MUTED,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(5.5),
        Inches(11),
        Inches(0.5),
        "CONFIDENTIAL — NOT LEGAL ADVICE",
        font_size=14,
        bold=True,
        color_hex=BRAND_PAPER,
    )

    add_text_box(
        slide,
        Inches(1),
        Inches(6.0),
        Inches(11),
        Inches(0.4),
        branding.legal_marking,
        font_size=10,
        color_hex=BRAND_ON_INK_MUTED,
    )

    add_speaker_notes(
        slide,
        "Cover slide. This report is AI-assisted and does not constitute legal advice.",
    )


def add_disclaimer_slide(prs, report: FTOReport, branding: BrandingConfig) -> None:
    """Slide 2: Disclaimer."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Important Disclaimer")
    disclaimer = branding.effective_disclaimer_text

    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(5.5),
        trunc(disclaimer, 1500),
        font_size=12,
        color_hex=BRAND_SECONDARY_TEXT,
    )

    add_speaker_notes(slide, "Read the key points of the disclaimer to the audience.")


def add_evidence_scope_slide(
    prs,
    report: FTOReport,
    *,
    options: ExportRenderOptions | None = None,
) -> None:
    """Slide 3: Evidence scope before risk conclusions."""
    from pptx.util import Inches

    from praviar_pipeline.rendering.export_options import default_export_options

    options = options or default_export_options()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Evidence Scope")

    summary = summarize_evidence_scope(report)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.35),
        Inches(5.5),
        Inches(0.72),
        summary.headline,
        font_size=24,
        bold=True,
        color_hex=BRAND_INK,
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.12),
        Inches(5.5),
        Inches(1.0),
        summary.posture,
        font_size=14,
        color_hex=BRAND_SECONDARY_TEXT,
    )

    scope_text = (
        "Successful sources\n"
        f"{format_source_list(summary.successful_sources, empty='No telemetry recorded')}\n\n"
        "Unavailable sources\n"
        f"{format_source_list(summary.unavailable_sources, empty='None recorded')}\n\n"
        "Skipped sources\n"
        f"{format_source_list(summary.skipped_sources, empty='None recorded')}"
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(3.35),
        Inches(5.5),
        Inches(2.65),
        scope_text,
        font_size=12,
        color_hex=BRAND_INK,
    )

    source_lines = []
    for entry in report.source_health.entries[:8]:
        source_lines.append(f"{entry.source}: {source_status_detail(entry)}")
    if len(report.source_health.entries) > 8:
        source_lines.append(f"+{len(report.source_health.entries) - 8} more sources")
    telemetry = "\n".join(source_lines) or "Source-health telemetry was not recorded."
    add_text_box(
        slide,
        Inches(6.8),
        Inches(1.45),
        Inches(5.7),
        Inches(3.55),
        telemetry,
        font_size=11,
        color_hex=BRAND_SECONDARY_TEXT,
    )
    add_text_box(
        slide,
        Inches(6.8),
        Inches(5.2),
        Inches(5.7),
        Inches(1.0),
        f"Confidence impact: {summary.confidence_impact}\n{summary.review_note}",
        font_size=12,
        bold=True,
        color_hex=BRAND_INK,
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(6.52),
        Inches(11.7),
        Inches(0.42),
        f"Audience: {options.audience_label} | Sections: {', '.join(options.section_labels)}",
        font_size=9,
        color_hex=BRAND_SECONDARY_TEXT,
    )

    add_speaker_notes(
        slide,
        f"Evidence scope: {summary.headline}. {summary.confidence_impact} {summary.review_note}",
    )


def add_executive_summary_slide(prs, report: FTOReport) -> None:
    """Slide 3: Executive summary with verdict and key metrics."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Executive Summary")

    governed_risk = governed_risk_level(report)
    risk = risk_display(governed_risk)
    risk_color = RISK_FILL.get(governed_risk, BRAND_SECONDARY_TEXT)

    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(4),
        Inches(0.8),
        f"Clearance Decision: {governed_decision_label(report)} ({risk} risk)",
        font_size=28,
        bold=True,
        color_hex=risk_color,
    )

    metrics = (
        f"Blocking Patents: {governed_blocking_count(report)}\n"
        f"Total Analyzed: {report.risk_summary.total_patents_analyzed}\n"
        f"Total Discovered: {report.total_patents_found}"
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.3),
        Inches(4),
        Inches(1.5),
        metrics,
        font_size=14,
        color_hex=BRAND_SECONDARY_TEXT,
    )

    summary = governed_executive_summary(report)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(4.2),
        Inches(11.5),
        Inches(2.8),
        summary,
        font_size=12,
        color_hex=BRAND_SECONDARY_TEXT,
    )

    add_speaker_notes(
        slide,
        f"Clearance decision is {governed_decision_label(report)} ({risk} risk). "
        f"{governed_blocking_count(report)} blocking patents identified "
        f"out of {report.risk_summary.total_patents_analyzed} analyzed.",
    )


def add_compound_slide(prs, report: FTOReport) -> None:
    """Slide 4: Compound overview."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, f"Compound: {report.compound.name}")

    c = report.compound
    props = (
        f"SMILES: {trunc(c.canonical_smiles, 80)}\n"
        f"InChIKey: {c.inchi_key}\n"
        f"Formula: {c.molecular_formula}"
        + (f"\nMW: {c.molecular_weight:.2f}" if c.molecular_weight else "")
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(6),
        Inches(3),
        props,
        font_size=12,
        color_hex=BRAND_INK,
    )

    if c.functional_groups:
        fg_text = "Functional Groups: " + ", ".join(c.functional_groups[:10])
        add_text_box(
            slide,
            Inches(0.8),
            Inches(4.5),
            Inches(11),
            Inches(0.5),
            fg_text,
            font_size=11,
            color_hex=BRAND_SECONDARY_TEXT,
        )

    add_speaker_notes(slide, f"Target compound: {c.name}. {c.molecular_formula}.")


def add_methodology_slide(prs, report: FTOReport) -> None:
    """Slide 5: Search methodology."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Search Methodology")

    summary = summarize_evidence_scope(report)
    if report.source_health.entries:
        sources = ", ".join(entry.source for entry in report.source_health.entries)
        source_label = "Configured Source Requests"
    else:
        sources = ", ".join(report.search_sources_used) or "Not recorded"
        source_label = "Recorded Source Names"
    text = (
        f"{source_label}: {sources}\n"
        f"Source Status: {summary.headline}\n\n"
        f"Total Patents Discovered: {report.total_patents_found:,}\n"
        f"After Triage: {report.patents_after_triage}\n"
        f"Analyzed in Detail: {len(report.patent_analyses)}\n\n"
        f"Execution Profile: {report.execution_profile}"
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(6),
        Inches(5),
        text,
        font_size=13,
        color_hex=BRAND_INK,
    )

    if report.source_health.entries:
        health_text = "Source Health:\n"
        for e in report.source_health.entries:
            health_text += f"  {source_status_label(e)} - {e.source}: {source_status_detail(e)}\n"
        add_text_box(
            slide,
            Inches(7),
            Inches(1.4),
            Inches(5.5),
            Inches(5),
            health_text,
            font_size=11,
            color_hex=BRAND_SECONDARY_TEXT,
        )

    add_speaker_notes(
        slide,
        f"Source telemetry: {summary.headline}. "
        f"The run found {report.total_patents_found} initial patents before triage.",
    )
