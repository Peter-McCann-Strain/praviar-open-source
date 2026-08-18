"""Shared layout helpers for PPTX report rendering."""

from __future__ import annotations

import io

from praviar_pipeline.rendering.design import BRAND_INK, BRAND_PAPER, hex_to_rgb


def trunc(text: str, max_len: int) -> str:
    """Truncate text with ellipsis if needed."""
    if len(text) <= max_len:
        return text
    return text[: max_len - 1] + "\u2026"


def pptx_rgb(hex_str: str):
    """Create a python-pptx RGBColor from hex string."""
    from pptx.dml.color import RGBColor

    r, g, b = hex_to_rgb(hex_str)
    return RGBColor(r, g, b)


def apply_font(
    run,
    font_name: str = "Calibri",
    size_pt: float = 14,
    bold: bool = False,
    italic: bool = False,
    color_hex: str = BRAND_INK,
) -> None:
    """Apply font styling to a text run."""
    from pptx.util import Pt

    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = pptx_rgb(color_hex)


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: float = 14,
    bold: bool = False,
    italic: bool = False,
    color_hex: str = BRAND_INK,
    alignment=None,
) -> None:
    """Add a text box to a slide."""
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    for run in p.runs:
        apply_font(run, size_pt=font_size, bold=bold, italic=italic, color_hex=color_hex)


def add_ink_title_bar(slide, title_text: str) -> None:
    """Add an Ink title bar across the top of a content slide."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = pptx_rgb(BRAND_INK)
    shape.line.fill.background()

    tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.6))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    apply_font(run, size_pt=24, bold=True, color_hex=BRAND_PAPER)


def add_speaker_notes(slide, notes: str) -> None:
    """Add speaker notes to a slide."""
    if not slide.has_notes_slide:
        _ = slide.notes_slide  # Create notes slide
    slide.notes_slide.notes_text_frame.text = notes


def add_chart_image(slide, chart_base64: str, left, top, width) -> None:
    """Embed a base64 PNG chart image on a slide."""
    import base64

    try:
        img_bytes = base64.b64decode(chart_base64)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, left, top, width=width)
    except Exception:
        import structlog

        structlog.get_logger().warning("pptx_chart_embed_failed")
