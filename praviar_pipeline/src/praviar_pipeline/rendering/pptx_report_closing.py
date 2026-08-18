"""Patent detail and closing slides for PPTX report rendering."""

from __future__ import annotations

from typing import TYPE_CHECKING

from praviar_pipeline.models.analysis import RiskLevel
from praviar_pipeline.rendering.design import (
    BRAND_INK,
    BRAND_ON_INK_MUTED,
    BRAND_PAPER,
    BRAND_SECONDARY_TEXT,
    RISK_FILL,
    risk_display,
    risk_sort_key,
)
from praviar_pipeline.rendering.pptx_report_layout import (
    add_ink_title_bar,
    add_speaker_notes,
    add_text_box,
    pptx_rgb,
    trunc,
)

if TYPE_CHECKING:
    from praviar_pipeline.models.analysis import PatentAnalysis
    from praviar_pipeline.models.report import FTOReport
    from praviar_pipeline.rendering.branding import BrandingConfig


def add_patent_deep_dive_slide(prs, analysis: PatentAnalysis, report: FTOReport) -> None:
    """Add 1 slide per HIGH-risk patent with key details."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    risk = risk_display(analysis.risk_level)
    risk_color = RISK_FILL.get(analysis.risk_level, BRAND_SECONDARY_TEXT)

    add_ink_title_bar(slide, f"{analysis.patent_id}: {trunc(analysis.title, 60)}")

    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.3),
        Inches(3),
        Inches(0.5),
        f"Risk: {risk}",
        font_size=18,
        bold=True,
        color_hex=risk_color,
    )

    metadata = (
        f"Assignee: {analysis.assignee}\n"
        f"Expiry: {analysis.expiry_date.isoformat() if analysis.expiry_date else 'Unknown'}\n"
        f"Claims analyzed: {len(analysis.claims_analyzed)}"
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.9),
        Inches(5),
        Inches(1.2),
        metadata,
        font_size=11,
        color_hex=BRAND_SECONDARY_TEXT,
    )

    narrative = report.patent_narratives.get(analysis.patent_id, analysis.risk_summary)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(3.3),
        Inches(7),
        Inches(3.5),
        trunc(narrative, 800),
        font_size=12,
        color_hex=BRAND_INK,
    )

    if analysis.claims_analyzed:
        claim_text = "Claim Summary:\n"
        for c in analysis.claims_analyzed[:5]:
            status = c.overall_status.value.upper().replace("_", " ")
            claim_text += f"  Claim {c.claim_number}: {status}\n"
        add_text_box(
            slide,
            Inches(8.5),
            Inches(1.3),
            Inches(4.5),
            Inches(3),
            claim_text,
            font_size=11,
            color_hex=BRAND_INK,
        )

    if analysis.design_around_suggestions:
        das_text = "Design-Around Options:\n"
        for sug in analysis.design_around_suggestions[:3]:
            das_text += f"  • {trunc(sug.suggestion, 100)} ({sug.feasibility})\n"
        add_text_box(
            slide,
            Inches(8.5),
            Inches(4.5),
            Inches(4.5),
            Inches(2.5),
            das_text,
            font_size=10,
            color_hex=BRAND_SECONDARY_TEXT,
        )

    add_speaker_notes(
        slide,
        f"Patent {analysis.patent_id} by {analysis.assignee}. Risk level: {risk}. "
        f"{len(analysis.claims_analyzed)} claims analyzed.",
    )


def add_moderate_summary_slide(prs, report: FTOReport) -> None:
    """Slide: Summary of moderate and low-risk patents."""
    from pptx.util import Inches

    moderate_low = [
        a
        for a in report.patent_analyses
        if a.risk_level in (RiskLevel.MEDIUM, RiskLevel.LOW, RiskLevel.CLEAR)
    ]
    if not moderate_low:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Other Patents Summary")

    text = ""
    for analysis in sorted(moderate_low, key=lambda x: risk_sort_key(x.risk_level))[:12]:
        risk = risk_display(analysis.risk_level)
        text += f"{analysis.patent_id} ({analysis.assignee}) — {risk}\n"

    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(5.5),
        text,
        font_size=12,
        color_hex=BRAND_INK,
    )

    add_speaker_notes(slide, f"Summary of {len(moderate_low)} moderate and low-risk patents.")


def add_recommendations_slide(prs, report: FTOReport) -> None:
    """Slide: Strategic recommendations."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Strategic Recommendations")

    if report.action_items:
        priority_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
        sorted_items = sorted(
            report.action_items,
            key=lambda x: priority_order.get(x.priority.value, 99),
        )
        text = ""
        for item in sorted_items[:8]:
            action = item.action_type.value.replace("_", " ").title()
            pri = item.priority.value.upper()
            text += f"[{pri}] {action}: {trunc(item.description, 100)}\n"
            if item.patent_ids:
                text += f"    Patents: {', '.join(item.patent_ids[:3])}\n"
            text += "\n"
    else:
        text = "No specific action items generated. Consult with patent counsel."

    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(5.5),
        text,
        font_size=13,
        color_hex=BRAND_INK,
    )

    add_speaker_notes(slide, "Key recommendations for next steps based on the analysis findings.")


def add_appendix_slide(prs, report: FTOReport, branding: BrandingConfig) -> None:
    """Final slide: Appendix / contact."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    version_label = (
        f"Report version {report.praviar_pipeline_version}"
        if branding.suppresses_praviar_branding
        else f"Praviar v{report.praviar_pipeline_version}"
    )
    notes = (
        "Thank you. This report was generated for counsel review. "
        "All findings should be reviewed by qualified patent counsel."
        if branding.suppresses_praviar_branding
        else "Thank you. This report was generated by Praviar. "
        "All findings should be reviewed by qualified patent counsel."
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = pptx_rgb(branding.primary_color)
    bg.line.fill.background()

    add_text_box(
        slide,
        Inches(1),
        Inches(2),
        Inches(11),
        Inches(1),
        branding.display_name,
        font_size=32,
        bold=True,
        color_hex=BRAND_PAPER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(3.5),
        Inches(11),
        Inches(0.5),
        version_label,
        font_size=14,
        color_hex=BRAND_ON_INK_MUTED,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(4.2),
        Inches(11),
        Inches(0.5),
        f"Generated: {report.generated_at.strftime('%Y-%m-%d %H:%M UTC')}",
        font_size=12,
        color_hex=BRAND_ON_INK_MUTED,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(5.5),
        Inches(11),
        Inches(0.5),
        "CONFIDENTIAL — NOT LEGAL ADVICE",
        font_size=12,
        bold=True,
        color_hex=BRAND_PAPER,
    )

    add_speaker_notes(slide, notes)
