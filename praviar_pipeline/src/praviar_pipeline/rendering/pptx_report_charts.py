"""Chart and matrix slides for PPTX report rendering."""

from __future__ import annotations

from collections import Counter
from typing import TYPE_CHECKING

from praviar_pipeline.rendering.design import (
    BRAND_INK,
    BRAND_PAPER,
    BRAND_SECONDARY_TEXT,
    RISK_BG,
    risk_display,
    risk_sort_key,
)
from praviar_pipeline.rendering.pptx_report_layout import (
    add_chart_image,
    add_ink_title_bar,
    add_speaker_notes,
    add_text_box,
    apply_font,
    pptx_rgb,
    trunc,
)

if TYPE_CHECKING:
    from praviar_pipeline.models.report import FTOReport


def add_funnel_slide(prs, report: FTOReport, charts: dict[str, str]) -> None:
    """Slide 6: Patent screening funnel chart."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Patent Screening Funnel")

    if "funnel" in charts:
        add_chart_image(slide, charts["funnel"], Inches(1.5), Inches(1.4), Inches(10))
    else:
        at = report.audit_trail
        text = (
            f"Discovered: {at.total_patents_discovered}\n"
            f"After Filters: {at.patents_after_hard_filter}\n"
            f"After Ranking: {at.patents_after_ranking}\n"
            f"After Triage: {at.patents_after_triage}\n"
            f"Analyzed: {at.patents_analyzed}"
        )
        add_text_box(
            slide,
            Inches(2),
            Inches(2),
            Inches(9),
            Inches(4),
            text,
            font_size=18,
            color_hex=BRAND_INK,
        )

    add_speaker_notes(slide, "This funnel shows how patents were progressively filtered.")


def add_risk_distribution_slide(prs, report: FTOReport, charts: dict[str, str]) -> None:
    """Slide 7: Risk distribution chart."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Risk Distribution")

    if "risk_distribution" in charts:
        add_chart_image(slide, charts["risk_distribution"], Inches(3), Inches(1.4), Inches(7))
    else:
        counts = Counter(risk_display(a.risk_level) for a in report.patent_analyses)
        text = "\n".join(f"{k}: {v}" for k, v in counts.most_common())
        add_text_box(
            slide,
            Inches(3),
            Inches(2),
            Inches(7),
            Inches(4),
            text,
            font_size=24,
            bold=True,
            color_hex=BRAND_INK,
        )

    add_speaker_notes(slide, "Distribution of patent risk levels across analyzed patents.")


def add_risk_matrix_slide(prs, report: FTOReport) -> None:
    """Slide 8: Risk matrix table."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Risk Assessment Matrix")

    if not report.patent_analyses:
        add_text_box(
            slide,
            Inches(3),
            Inches(3),
            Inches(7),
            Inches(1),
            "No patents analyzed",
            font_size=18,
            color_hex=BRAND_SECONDARY_TEXT,
        )
        return

    sorted_analyses = sorted(report.patent_analyses, key=lambda a: risk_sort_key(a.risk_level))
    display = sorted_analyses[:15]

    table_shape = slide.shapes.add_table(
        len(display) + 1,
        5,
        Inches(0.4),
        Inches(1.3),
        Inches(12.5),
        Inches(5.8),
    )
    table = table_shape.table

    headers = ["Patent ID", "Title", "Assignee", "Risk", "Expiry"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                apply_font(run, size_pt=10, bold=True, color_hex=BRAND_PAPER)
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx_rgb(BRAND_SECONDARY_TEXT)

    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(1.5)
    table.columns[4].width = Inches(2.0)

    for row_idx, a in enumerate(display, 1):
        table.cell(row_idx, 0).text = a.patent_id
        table.cell(row_idx, 1).text = trunc(a.title, 50)
        table.cell(row_idx, 2).text = trunc(a.assignee, 30)
        table.cell(row_idx, 3).text = risk_display(a.risk_level)
        table.cell(row_idx, 4).text = a.expiry_date.isoformat() if a.expiry_date else "N/A"

        risk_bg = RISK_BG.get(a.risk_level, BRAND_PAPER)
        table.cell(row_idx, 3).fill.solid()
        table.cell(row_idx, 3).fill.fore_color.rgb = pptx_rgb(risk_bg)

        for col_idx in range(5):
            for p in table.cell(row_idx, col_idx).text_frame.paragraphs:
                for run in p.runs:
                    apply_font(run, size_pt=9, color_hex=BRAND_INK)

    if len(sorted_analyses) > 15:
        add_text_box(
            slide,
            Inches(0.8),
            Inches(7.0),
            Inches(5),
            Inches(0.3),
            f"Showing {15} of {len(sorted_analyses)} patents",
            font_size=8,
            italic=True,
            color_hex=BRAND_SECONDARY_TEXT,
        )

    add_speaker_notes(slide, f"Risk matrix showing {len(display)} patents sorted by risk level.")


def add_timeline_slide(prs, report: FTOReport, charts: dict[str, str]) -> None:
    """Slide 9: Patent timeline Gantt chart."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ink_title_bar(slide, "Patent Expiry Timeline")

    if "timeline" in charts:
        add_chart_image(slide, charts["timeline"], Inches(1), Inches(1.4), Inches(11))
    else:
        patents_with_expiry = [
            (analysis, expiry_date)
            for analysis in report.patent_analyses
            if (expiry_date := analysis.expiry_date) is not None
        ]
        patents_with_expiry.sort(key=lambda item: item[1])
        text = "\n".join(
            (
                f"{analysis.patent_id}: expires {expiry_date.isoformat()} "
                f"({risk_display(analysis.risk_level)})"
            )
            for analysis, expiry_date in patents_with_expiry[:12]
        )
        add_text_box(
            slide,
            Inches(1),
            Inches(1.5),
            Inches(11),
            Inches(5.5),
            text or "No expiry dates available",
            font_size=12,
            color_hex=BRAND_INK,
        )

    add_speaker_notes(
        slide,
        "Timeline showing when each patent expires. Opportunities open as patents expire.",
    )
