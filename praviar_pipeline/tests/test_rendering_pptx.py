"""Tests for praviar_pipeline.rendering.pptx_report -- PPTX export for FTO reports."""
# ruff: noqa: E402

from __future__ import annotations

import base64
import io
from datetime import UTC, date, datetime

import pytest

pptx = pytest.importorskip("pptx")

from praviar_pipeline.models.analysis import (
    ClaimAnalysis,
    ClaimElement,
    ElementStatus,
    PatentAnalysis,
    RiskLevel,
)
from praviar_pipeline.models.audit import PipelineAuditTrail, SearchFunnelEntry, StepTiming
from praviar_pipeline.models.compound import ResolvedCompound
from praviar_pipeline.models.report import (
    FTOReport,
    RiskSummary,
    SourceHealth,
    SourceHealthEntry,
    SourceStatus,
)
from praviar_pipeline.models.verification import VerificationCheck, VerificationResult
from praviar_pipeline.rendering.branding import BrandingConfig
from praviar_pipeline.rendering.export_options import ExportRenderOptions
from praviar_pipeline.rendering.pptx_report import render_pptx

_PNG_1X1 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def minimal_compound() -> ResolvedCompound:
    return ResolvedCompound(
        name="Aspirin",
        canonical_smiles="CC(=O)Oc1ccccc1C(=O)O",
        inchi="InChI=1S/C9H8O4/c1-6(10)13-8-5-3-2-4-7(8)9(11)12/h2-5H,1H3,(H,11,12)",
        inchi_key="BSYNRYMUTXBXSQ-UHFFFAOYSA-N",
        molecular_formula="C9H8O4",
        molecular_weight=180.16,
        pubchem_cid=2244,
        original_input="aspirin",
        input_type="name",
    )


@pytest.fixture
def medium_analysis() -> PatentAnalysis:
    return PatentAnalysis(
        patent_id="US7851188B2",
        title="Methods for producing aspirin derivatives",
        assignee="PharmaCorp Inc.",
        expiry_date=date(2028, 3, 15),
        claims_analyzed=[
            ClaimAnalysis(
                claim_number=1,
                claim_type="independent",
                elements=[
                    ClaimElement(
                        element_number=1,
                        element_text="A method for producing acetylsalicylic acid",
                        status=ElementStatus.MET,
                        reasoning="Target compound is acetylsalicylic acid (aspirin)",
                        confidence=0.95,
                    ),
                    ClaimElement(
                        element_number=2,
                        element_text="comprising reacting salicylic acid with acetic anhydride",
                        status=ElementStatus.NOT_MET,
                        reasoning="Different synthesis route used",
                        confidence=0.9,
                    ),
                ],
                overall_status=ElementStatus.NOT_MET,
                overall_confidence=0.9,
            ),
        ],
        risk_level=RiskLevel.MEDIUM,
        risk_summary="Moderate risk -- one claim element not met",
    )


@pytest.fixture
def high_risk_analysis() -> PatentAnalysis:
    return PatentAnalysis(
        patent_id="US8888888B2",
        title="Aspirin formulation process",
        assignee="GreenChem Corp",
        expiry_date=date(2035, 1, 1),
        claims_analyzed=[
            ClaimAnalysis(
                claim_number=1,
                claim_type="independent",
                elements=[
                    ClaimElement(
                        element_number=1,
                        element_text="producing acetylsalicylic acid",
                        status=ElementStatus.MET,
                        reasoning="Exact match",
                        confidence=0.99,
                    ),
                ],
                overall_status=ElementStatus.MET,
                overall_confidence=0.99,
            ),
        ],
        risk_level=RiskLevel.HIGH,
        risk_summary="All claim elements met -- direct infringement risk",
    )


@pytest.fixture
def rendering_report(
    minimal_compound: ResolvedCompound,
    medium_analysis: PatentAnalysis,
    high_risk_analysis: PatentAnalysis,
) -> FTOReport:
    """FTOReport with all fields populated for rendering tests."""
    now = datetime.now(UTC)
    return FTOReport(
        report_id="pptx-test-001",
        generated_at=now,
        compound=minimal_compound,
        risk_summary=RiskSummary(
            overall_risk=RiskLevel.HIGH,
            blocking_patents_count=2,
            total_patents_analyzed=5,
            key_risks=[
                "US7851188B2: medium risk",
                "US8888888B2: blocks composition",
            ],
            executive_summary="Analysis reveals 2 blocking patents for Aspirin.",
        ),
        patent_analyses=[medium_analysis, high_risk_analysis],
        verification=VerificationResult(
            checks=[
                VerificationCheck(
                    check_name="citation_grounding",
                    passed=True,
                    details="All patent IDs found in search results",
                ),
            ],
            all_citations_valid=True,
            all_claims_grounded=True,
            all_entities_valid=True,
            dates_consistent=True,
            risk_levels_justified=True,
        ),
        total_patents_found=50,
        patents_after_triage=2,
        search_sources_used=["pubchem", "bigquery"],
        source_health=SourceHealth(
            entries=[
                SourceHealthEntry(
                    source="pubchem",
                    status=SourceStatus.OK,
                    patent_count=42,
                ),
                SourceHealthEntry(
                    source="bigquery",
                    status=SourceStatus.FAILED,
                    patent_count=0,
                    error_message="HTTP 504 timeout from BigQuery",
                ),
                SourceHealthEntry(
                    source="lens",
                    status=SourceStatus.NOT_CONFIGURED,
                    patent_count=0,
                    error_message="API key missing",
                ),
            ]
        ),
        audit_trail=PipelineAuditTrail(
            search_funnel=[
                SearchFunnelEntry(patent_id="US7851188B2", included_in_triage=True),
            ],
            timing_data=[
                StepTiming(
                    step_name="step2_search",
                    started_at=now,
                    completed_at=now,
                    duration_seconds=3.5,
                    items_processed=100,
                    items_output=50,
                ),
            ],
            total_patents_discovered=100,
        ),
        patent_narratives={
            "US7851188B2": "This patent covers aspirin derivative production methods.",
        },
        llm_models_used={
            "triage": "claude-haiku-4-5-20251001",
            "analysis": "claude-sonnet-4-20250514",
        },
    )


@pytest.fixture
def empty_analyses_report(minimal_compound: ResolvedCompound) -> FTOReport:
    """FTOReport with no patent_analyses (edge case)."""
    return FTOReport(
        report_id="pptx-test-empty",
        compound=minimal_compound,
        risk_summary=RiskSummary(
            overall_risk=RiskLevel.CLEAR,
            blocking_patents_count=0,
            total_patents_analyzed=0,
            key_risks=[],
            executive_summary="No blocking patents identified.",
        ),
        patent_analyses=[],
        total_patents_found=0,
        patents_after_triage=0,
        search_sources_used=["pubchem"],
    )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _all_slide_text(prs) -> str:
    """Extract all text content from all slides in a Presentation."""
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts)


def _all_slide_notes_text(prs) -> str:
    """Extract speaker notes from all slides in a Presentation."""
    parts = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text:
                parts.append(notes_text)
    return "\n".join(parts)


def _slide_text(slide) -> str:
    """Extract all text from a single slide."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                parts.append(para.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


def _write_png_logo(path) -> None:
    path.write_bytes(base64.b64decode(_PNG_1X1))


def _slide_has_table(slide) -> bool:
    """Check whether a slide contains at least one table."""
    return any(shape.has_table for shape in slide.shapes)


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestPptxRendering:
    """Tests for the PPTX report renderer."""

    def test_returns_bytes(self, rendering_report: FTOReport):
        """render_pptx should return bytes."""
        result = render_pptx(rendering_report)
        assert isinstance(result, bytes)
        assert len(result) > 0

    def test_zip_magic_bytes(self, rendering_report: FTOReport):
        """PPTX files are ZIP archives starting with PK header."""
        result = render_pptx(rendering_report)
        assert result[:4] == b"PK\x03\x04"

    def test_loadable_as_presentation(self, rendering_report: FTOReport):
        """The returned bytes should load as a valid python-pptx Presentation."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        assert prs.slides is not None

    def test_executive_verdict_uses_clearance_decision(self, rendering_report: FTOReport):
        from pptx import Presentation

        rendering_report.risk_summary.overall_risk = RiskLevel.CLEAR
        result = render_pptx(rendering_report)
        text = _all_slide_text(Presentation(io.BytesIO(result)))

        assert "Clearance Decision: UNCLEAR (MODERATE risk)" in text
        assert "Overall Risk: Clear" not in text

    def test_executive_scope_omits_analysis_slides(self, rendering_report: FTOReport):
        """Executive-only PPTX exports keep caveats but omit technical slides."""
        from pptx import Presentation

        options = ExportRenderOptions.from_values(
            ["executive_summary"],
            audience="investor",
        )
        result = render_pptx(rendering_report, options=options)
        prs = Presentation(io.BytesIO(result))
        all_text = _all_slide_text(prs)

        assert "Important Disclaimer" in all_text
        assert "Evidence Scope" in all_text
        assert "Audience: Investor Pack" in all_text
        assert "Executive Summary" in all_text
        assert "Strategic Recommendations" in all_text
        assert "Risk Matrix" not in all_text
        assert "Search Methodology" not in all_text

    def test_minimum_slide_count(self, rendering_report: FTOReport):
        """Presentation should have at least 10 slides for a report with analyses."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        # Cover + Disclaimer + Executive Summary + Compound + Methodology +
        # Funnel + Risk Distribution + Risk Matrix + Timeline + deep dives +
        # Moderate summary + Recommendations + Appendix = 13+ slides
        assert len(prs.slides) >= 10

    def test_cover_slide_title(self, rendering_report: FTOReport):
        """Slide 1 (cover) should contain the FTO analysis title."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        cover_text = _slide_text(prs.slides[0])
        assert "Freedom-to-Operate" in cover_text

    def test_cover_slide_compound_name(self, rendering_report: FTOReport):
        """Cover slide should contain the compound name."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        cover_text = _slide_text(prs.slides[0])
        assert "Aspirin" in cover_text

    def test_cover_slide_report_id(self, rendering_report: FTOReport):
        """Cover slide should contain the report ID."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        cover_text = _slide_text(prs.slides[0])
        assert "pptx-test-001" in cover_text

    def test_evidence_scope_slide_precedes_executive_summary(
        self,
        rendering_report: FTOReport,
    ):
        """Slide 3 should carry source-health posture before risk conclusions."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        evidence_text = _slide_text(prs.slides[2])
        assert "Evidence Scope" in evidence_text
        assert "1 of 3 configured sources completed" in evidence_text
        assert "bigquery" in evidence_text
        assert "Confidence impact" in evidence_text

    def test_executive_summary_slide(self, rendering_report: FTOReport):
        """Slide 4 (exec summary) should contain risk level text."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        # Slide 4 is index 3 (0=cover, 1=disclaimer, 2=evidence, 3=exec).
        exec_text = _slide_text(prs.slides[3])
        assert "Clearance Decision: UNCLEAR (MODERATE risk)" in exec_text

    def test_executive_summary_key_risks(self, rendering_report: FTOReport):
        """Executive summary must not reintroduce drifting risk prose."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        exec_text = _slide_text(prs.slides[3])
        assert "Clearance decision: UNCLEAR." in exec_text
        assert "blocks composition" not in exec_text

    def test_risk_matrix_slide_has_table(self, rendering_report: FTOReport):
        """Risk matrix slide should contain a table."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        # Risk matrix is slide 9 (index 8) after the Evidence Scope slide.
        risk_matrix_slide = prs.slides[8]
        assert _slide_has_table(risk_matrix_slide)

    def test_risk_matrix_has_patent_data(self, rendering_report: FTOReport):
        """Risk matrix table should contain patent IDs."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        all_text = _all_slide_text(prs)
        assert "US7851188B2" in all_text
        assert "US8888888B2" in all_text

    def test_deep_dive_slides_for_high_risk(self, rendering_report: FTOReport):
        """High-risk patents should get their own deep-dive slide."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        all_text = _all_slide_text(prs)
        assert "US8888888B2" in all_text
        assert "GreenChem Corp" in all_text

    def test_all_text_contains_compound(self, rendering_report: FTOReport):
        """Compound details should appear somewhere in the presentation."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        all_text = _all_slide_text(prs)
        assert "Aspirin" in all_text
        assert "CC(=O)Oc1ccccc1C(=O)O" in all_text

    def test_confidentiality_marking(self, rendering_report: FTOReport):
        """Confidentiality marking should appear on cover slide."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        cover_text = _slide_text(prs.slides[0])
        assert "CONFIDENTIAL" in cover_text


class TestPptxSpeakerNotes:
    """Tests for speaker notes in PPTX slides."""

    def test_cover_slide_has_notes(self, rendering_report: FTOReport):
        """Cover slide should have speaker notes."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[0]
        assert slide.has_notes_slide
        notes_text = slide.notes_slide.notes_text_frame.text
        assert len(notes_text) > 0

    def test_disclaimer_slide_has_notes(self, rendering_report: FTOReport):
        """Disclaimer slide should have speaker notes."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[1]
        assert slide.has_notes_slide
        notes_text = slide.notes_slide.notes_text_frame.text
        assert len(notes_text) > 0

    def test_exec_summary_notes_contain_risk(self, rendering_report: FTOReport):
        """Executive summary notes should mention the risk level."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[3]
        assert slide.has_notes_slide
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Clearance decision is UNCLEAR (MODERATE risk)" in notes_text

    def test_multiple_slides_have_notes(self, rendering_report: FTOReport):
        """Most slides should have speaker notes."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        slides_with_notes = 0
        for slide in prs.slides:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slides_with_notes += 1
        # At least half the slides should have notes
        assert slides_with_notes >= len(prs.slides) // 2


class TestPptxBranding:
    """Tests for PPTX rendering with BrandingConfig."""

    def test_with_privilege_header(self, rendering_report: FTOReport):
        """Organization branding cannot self-assert attorney privilege."""
        from pptx import Presentation

        branding = BrandingConfig(
            privilege_header="PRIVILEGED AND CONFIDENTIAL -- ATTORNEY WORK PRODUCT",
        )
        result = render_pptx(rendering_report, branding=branding)
        prs = Presentation(io.BytesIO(result))
        cover_text = _slide_text(prs.slides[0])
        assert "PRIVILEGED AND CONFIDENTIAL" not in cover_text
        assert "CONFIDENTIAL DRAFT" in cover_text

    def test_with_firm_name(self, rendering_report: FTOReport):
        """Firm name should appear in the appendix/closing slide."""
        from pptx import Presentation

        branding = BrandingConfig(firm_name="Baker & McKenzie LLP")
        result = render_pptx(rendering_report, branding=branding)
        prs = Presentation(io.BytesIO(result))
        all_text = _all_slide_text(prs)
        assert "Baker & McKenzie LLP" in all_text

    def test_default_branding_produces_valid_pptx(self, rendering_report: FTOReport):
        """Rendering with default branding should produce valid PPTX."""
        result = render_pptx(rendering_report, branding=None)
        assert isinstance(result, bytes)
        assert result[:4] == b"PK\x03\x04"

    def test_default_branding_includes_praviar_mark(self, rendering_report: FTOReport):
        """Default PPTX cover should include the Praviar mark image."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        result = render_pptx(rendering_report, branding=None)
        prs = Presentation(io.BytesIO(result))
        cover = prs.slides[0]
        assert "Praviar" in _slide_text(cover)
        assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in cover.shapes)

    def test_hidden_branding_removes_praviar_from_slides_and_notes(
        self,
        rendering_report: FTOReport,
    ):
        """Full white-label mode should not leak Praviar in visible text or notes."""
        from pptx import Presentation

        branding = BrandingConfig(hide_praviar_pipeline_branding=True)
        result = render_pptx(rendering_report, branding=branding)
        prs = Presentation(io.BytesIO(result))

        assert "Praviar" not in _all_slide_text(prs)
        assert "Praviar" not in _all_slide_notes_text(prs)
        assert "Report version" in _all_slide_text(prs)

    def test_custom_disclaimer_reaches_disclaimer_slide(
        self,
        rendering_report: FTOReport,
    ):
        """PPTX should honor BrandingConfig.disclaimer_text."""
        from pptx import Presentation

        branding = BrandingConfig(disclaimer_text="Acme custom export disclaimer.")
        result = render_pptx(rendering_report, branding=branding)
        prs = Presentation(io.BytesIO(result))

        disclaimer_text = _all_slide_text(prs)
        assert "Acme custom export disclaimer." in disclaimer_text
        assert "does NOT constitute legal advice" in disclaimer_text

    def test_custom_png_logo_uses_firm_identity_without_praviar_text(
        self,
        rendering_report: FTOReport,
        tmp_path,
    ):
        """Custom logo decks should not pair the firm mark with Praviar copy."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        logo = tmp_path / "firm-logo.png"
        _write_png_logo(logo)
        branding = BrandingConfig(logo_path=str(logo), firm_name="Acme Counsel")

        result = render_pptx(rendering_report, branding=branding)
        prs = Presentation(io.BytesIO(result))
        cover = prs.slides[0]

        assert "Acme Counsel" in _slide_text(cover)
        assert "Praviar" not in _all_slide_text(prs)
        assert "Praviar" not in _all_slide_notes_text(prs)
        assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in cover.shapes)

    def test_custom_svg_logo_fails_with_clear_pptx_error(
        self,
        rendering_report: FTOReport,
        tmp_path,
    ):
        """PPTX should reject unsupported SVG logos before python-pptx crashes."""
        logo = tmp_path / "firm-logo.svg"
        logo.write_text("<svg xmlns='http://www.w3.org/2000/svg' />", encoding="utf-8")
        branding = BrandingConfig(logo_path=str(logo), firm_name="Acme Counsel")

        with pytest.raises(
            RuntimeError,
            match=r"PPTX does not support branding logo format '\.svg'",
        ):
            render_pptx(rendering_report, branding=branding)

    def test_missing_custom_logo_fails_closed(
        self,
        rendering_report: FTOReport,
        tmp_path,
    ):
        """Missing logos should not silently produce an unbranded deck."""
        branding = BrandingConfig(
            logo_path=str(tmp_path / "missing-logo.png"),
            firm_name="Acme Counsel",
        )

        with pytest.raises(RuntimeError, match="Branding logo not found for PPTX"):
            render_pptx(rendering_report, branding=branding)


class TestPptxEdgeCases:
    """Tests for edge cases in PPTX rendering."""

    def test_empty_patent_analyses(self, empty_analyses_report: FTOReport):
        """Should produce valid PPTX even with no patent analyses."""
        result = render_pptx(empty_analyses_report)
        assert isinstance(result, bytes)
        assert result[:4] == b"PK\x03\x04"

    def test_empty_analyses_loadable(self, empty_analyses_report: FTOReport):
        """Empty analyses report should be loadable as a Presentation."""
        from pptx import Presentation

        result = render_pptx(empty_analyses_report)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) > 0

    def test_empty_analyses_risk_matrix_no_crash(self, empty_analyses_report: FTOReport):
        """Risk matrix slide should handle empty analyses gracefully."""
        from pptx import Presentation

        result = render_pptx(empty_analyses_report)
        prs = Presentation(io.BytesIO(result))
        # Should still have slides, no exception thrown
        all_text = _all_slide_text(prs)
        assert "No patents analyzed" in all_text or "CLEAR" in all_text

    def test_slide_width_is_widescreen(self, rendering_report: FTOReport):
        """Presentation should use widescreen (16:9 equivalent) dimensions."""
        from pptx import Presentation
        from pptx.util import Inches

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        # 13.333 inches width (set in render_pptx)
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)

    def test_methodology_slide_sources(self, rendering_report: FTOReport):
        """Methodology slide should list configured source telemetry."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        all_text = _all_slide_text(prs)
        assert "Configured Source Requests" in all_text
        assert "Databases Searched" not in all_text
        assert "pubchem" in all_text
        assert "bigquery" in all_text
        assert "lens" in all_text
        assert "Source Status" in all_text

    def test_provider_query_credentials_are_not_rendered(self, rendering_report: FTOReport):
        from pptx import Presentation

        rendering_report.source_health.entries[
            1
        ].error_message = "401 https://api.openalex.org/works?api_key=SUPERSECRET"

        all_text = _all_slide_text(Presentation(io.BytesIO(render_pptx(rendering_report))))

        assert "Provider request failed" in all_text
        assert "SUPERSECRET" not in all_text
        assert "api_key=" not in all_text

    def test_disclaimer_slide_content(self, rendering_report: FTOReport):
        """Disclaimer slide should contain key disclaimer text."""
        from pptx import Presentation

        result = render_pptx(rendering_report)
        prs = Presentation(io.BytesIO(result))
        disclaimer_text = _slide_text(prs.slides[1])
        assert "NOT constitute" in disclaimer_text or "does NOT" in disclaimer_text
